from pptx.dml.color import RGBColor
from pptx.util import Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN


def draw_rounder_rectangle(slide,left,top,width,height,color):
    
    rect = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left,
        top,
        width,
        height
    )
    
    rect.fill.solid()
    rect.fill.fore_color.rgb = color  
    rect.line.color.rgb = color
    
    left_circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        left - height/2,
        top,
        height,
        height
    )
    left_circle.fill.solid()
    left_circle.fill.fore_color.rgb = color 
    left_circle.line.color.rgb = color
    
    right_circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        left + width - height/2,
        top,
        height,
        height
    )
    right_circle.fill.solid()
    right_circle.fill.fore_color.rgb = color
    right_circle.line.color.rgb = color   

def add_textbox(slide,data,left,top,width,height,text_color):
    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    text_frame.text = data 
    paragraph = text_frame.paragraphs[0]
    run = paragraph.runs[0]
    run.font.name = 'Montserrat'
    run.font.size = Pt(15)
    run.font.color.rgb = text_color
    paragraph.alignment = PP_ALIGN.CENTER 


def draw_rounded_tf(slide, data, font_size, placeholder_index):
    placeholder_frame = slide.placeholders[placeholder_index]

    placeholder = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        placeholder_frame.left,
        placeholder_frame.top,
        placeholder_frame.width,
        placeholder_frame.height
    )

    frame = placeholder.text_frame
    frame.text = f"{data}"
    slide.shapes._spTree.remove(placeholder_frame._element)
    
    paragraph = frame.paragraphs[0]
    run = paragraph.runs[0]
    run.font.name = 'Montserrat'
    run.font.size = Pt(font_size)
    run.font.color.rgb = RGBColor(0, 0, 0)
    paragraph.alignment = PP_ALIGN.CENTER

    placeholder.fill.background()
    placeholder.line.fill.background()