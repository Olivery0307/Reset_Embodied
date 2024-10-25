from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from add_shapes_textbox import draw_rounder_rectangle, add_textbox
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN


def add_scorebar_border(slide,placeholder_idx):
    placeholder = slide.placeholders[placeholder_idx]
    
    left = placeholder.left
    top = placeholder.top
    width = placeholder.width
    height = placeholder.height
    color = RGBColor(0, 0, 0)  # Black 
    
    initial_shape_count = len(slide.shapes)
    draw_rounder_rectangle(slide,left, top, width, height, color)
    
    for i in range(initial_shape_count, len(slide.shapes)):
        shape = slide.shapes[i]
        shape.fill.background()  # Set fill to transparent
        shape.line.color.rgb = RGBColor(0, 0, 0)
        shape.line.width = Pt(3.8)  # Adjust this value to change border thickness
    draw_rounder_rectangle(slide, left, top, width, height, color = RGBColor(255,255,255))


def add_score_bar(slide, color, score, placeholder_idx):
    
    placeholder = slide.placeholders[placeholder_idx]
    
    left = placeholder.left
    top = placeholder.top
    width = placeholder.width
    height = placeholder.height

    # Calculate the score width
    score_width = width * score / 100
    draw_rounder_rectangle(slide,left,top,score_width,height,color)
    slide.shapes._spTree.remove(placeholder._element)
    
    #text
    if score == 100.0:
        draw_rounder_rectangle(slide,left + score_width - Inches(0.215),top+Inches(0.02),Inches(0.22),height-Inches(0.04),RGBColor(0xFF, 0xFF, 0xFF))
        add_textbox(slide,str(score),left+ score_width - Inches(0.215), top+Inches(0.02), Inches(0.22), height-Inches(0.04),RGBColor(0, 0, 0))  
        
    else:
        
        circle_size = height-Inches(0.02)
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            left + score_width - Inches(0.21),
            top + Inches(0.01),
            circle_size,
            circle_size)
        circle.fill.solid()
        circle.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        circle.line.color.rgb = RGBColor(0xFF, 0xFF, 0xFF) 
        
        textbox_size = circle_size * 1.5
        
        textbox = slide.shapes.add_textbox(
            left + score_width - Inches(0.13) - (textbox_size - circle_size) / 2,
            top+Inches(0.01),
            textbox_size - Inches(0.12),
            textbox_size)
        
        text_frame = textbox.text_frame
        text_frame.text = f"{score}"
        paragraph = text_frame.paragraphs[0]
        run = paragraph.runs[0]
        run.font.name = 'Montserrat' 
        run.font.size = Pt(15) 
        run.font.color.rgb = RGBColor(0, 0, 0) 
        paragraph.alignment = PP_ALIGN.CENTER    

def add_carbon_footprint(slide,data,idx):
    
    def format_number(value):
        try:
            num = float(value)
            int_part, dec_part = str(num).split('.')
            
            if len(dec_part) >= 2:
                return f"{int_part}.{dec_part[:2]}"
            else:
                return f"{int_part}.{dec_part}"
        except ValueError:
            return value
        
    placeholder = slide.placeholders[idx]
    left = placeholder.left
    top = placeholder.top
    width = placeholder.width
    height = placeholder.height
    
    initial_shape_count = len(slide.shapes)
    color = RGBColor(0, 0, 0)
    draw_rounder_rectangle(slide,left+Inches(0.16), top+Inches(0.01), width-Inches(0.26), height-Inches(0.01),color)
    
    for i in range(initial_shape_count, len(slide.shapes)):
        shape = slide.shapes[i]
        shape.fill.background()
        shape.line.color.rgb = RGBColor(0, 0, 0)
        shape.line.width = Pt(3)
    draw_rounder_rectangle(slide,left+Inches(0.16), top+Inches(0.01), width-Inches(0.26), height-Inches(0.01), color = RGBColor(255,255,255))
    
    draw_rounder_rectangle(slide,left+Inches(0.16), top+Inches(0.01), width-Inches(0.26),height-Inches(0.01),RGBColor(215,168,65))
    formatted_data = format_number(data)
    add_textbox(slide,f"{formatted_data}",left+Inches(0.16), top+Inches(0.01), width-Inches(0.26),height-Inches(0.01),RGBColor(0xFF, 0xFF, 0xFF))
    slide.shapes._spTree.remove(placeholder._element)


def add_donut_chart(slide,row,placeholder_idx=33):

    co2_value = float(row['RESET SCORE'].rstrip("%")) * 0.33
    circulation_value = ((float(row['RC'].rstrip("%")) + float(row['R'].rstrip("%"))) / 2) * 0.33
    health_value = 23
    remaining_value = 100 - (co2_value + circulation_value + health_value)
    
    chart_data = CategoryChartData()
    chart_data.categories = ['CO2','Circulation','Health','Remaining']
    chart_data.add_series('Series 1', (co2_value,circulation_value,health_value,remaining_value))

    placeholder = slide.placeholders[placeholder_idx]
    chart = placeholder.insert_chart(XL_CHART_TYPE.DOUGHNUT, chart_data).chart
    chart.has_legend = False
    chart.plots[0].has_data_labels = False
    chart.chart_title.text_frame.clear()
    
    colors = [
        RGBColor(0xD4, 0xA5, 0x23), # Gold
        RGBColor(187,198,200),  # Light gray
        RGBColor(193,163,131),  # Light brown
        RGBColor(235,235,235),  # White
    ]

    for i, point in enumerate(chart.series[0].points):
        point.format.fill.solid()
        point.format.fill.fore_color.rgb = colors[i]
        
    doughnut_element = chart.plots[0]._element
    holeSizes = doughnut_element.xpath('./c:holeSize')
    holeSize = holeSizes[0]
    holeSize.set('val', '75')